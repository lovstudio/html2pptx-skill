#!/usr/bin/env python3
"""
html2pptx — render an HTML file into a PPTX deck (one slide per page).

Pipeline:
  HTML → Playwright (Chromium headless) → per-slide PNG → python-pptx

Page-splitting strategies (auto-detected, override with --split):
  1. selector  — every element matching --selector (default: ".slide, section.slide")
                  is rendered as its own slide.
  2. pagebreak — split a single full-page render at CSS `page-break-after: always`
                  / `break-after: page` boundaries (handled by setting page size in
                  print emulation and walking the printed pages).
  3. single    — render the entire HTML as ONE slide (no splitting).

The default (`auto`) prefers `selector` if matching nodes exist, else `single`.

Examples:
  python3 html2pptx.py --input deck.html --output deck.pptx
  python3 html2pptx.py --input deck.html --selector ".slide" --size 1920x1080
  python3 html2pptx.py --input deck.html --split single --output single.pptx
  python3 html2pptx.py --input deck.html --scale 2 --output retina.pptx
"""

from __future__ import annotations

import argparse
import asyncio
import sys
import tempfile
from pathlib import Path
from typing import List, Tuple

DEFAULT_SELECTOR = ".slide, section.slide, [data-slide]"


def parse_size(s: str) -> Tuple[int, int]:
    try:
        w, h = s.lower().replace("×", "x").split("x")
        return int(w), int(h)
    except Exception:
        raise argparse.ArgumentTypeError(f"--size must be WxH (e.g. 1920x1080), got {s!r}")


async def render_via_selector(
    html_path: Path,
    out_dir: Path,
    selector: str,
    width: int,
    height: int,
    scale: float,
) -> List[Path]:
    """Render every element matching *selector* into its own PNG."""
    from playwright.async_api import async_playwright

    pngs: List[Path] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": width, "height": height},
            device_scale_factor=scale,
        )
        page = await ctx.new_page()
        await page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        # Wait extra tick for webfonts
        try:
            await page.evaluate("document.fonts && document.fonts.ready")
        except Exception:
            pass

        handles = await page.query_selector_all(selector)
        if not handles:
            await browser.close()
            return []

        for i, h in enumerate(handles, 1):
            # Force the element to be exactly slide-sized so screenshots line up.
            await h.evaluate(
                """(el, [w, h]) => {
                    el.style.width = w + 'px';
                    el.style.height = h + 'px';
                    el.style.boxSizing = 'border-box';
                    el.scrollIntoView({block: 'start', inline: 'start'});
                }""",
                [width, height],
            )
            png = out_dir / f"slide-{i:03d}.png"
            await h.screenshot(path=str(png), type="png")
            pngs.append(png)

        await browser.close()
    return pngs


async def render_single(
    html_path: Path,
    out_dir: Path,
    width: int,
    height: int,
    scale: float,
) -> List[Path]:
    """Render the whole page as a single PNG, viewport-sized."""
    from playwright.async_api import async_playwright

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": width, "height": height},
            device_scale_factor=scale,
        )
        page = await ctx.new_page()
        await page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        try:
            await page.evaluate("document.fonts && document.fonts.ready")
        except Exception:
            pass
        png = out_dir / "slide-001.png"
        await page.screenshot(path=str(png), type="png", full_page=False)
        await browser.close()
    return [png]


def assemble_pptx(pngs: List[Path], out_path: Path, width: int, height: int) -> None:
    """Pack PNGs into a 16:9 (or whatever) PPTX, one slide per image."""
    from pptx import Presentation
    from pptx.util import Emu

    # 1 inch = 914400 EMU. Translate pixel dims to EMU at 96 dpi.
    px_to_emu = 914400 / 96
    prs = Presentation()
    prs.slide_width = Emu(int(width * px_to_emu))
    prs.slide_height = Emu(int(height * px_to_emu))
    blank = prs.slide_layouts[6]  # blank layout

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


async def detect_split_mode(html_path: Path, selector: str) -> str:
    """Peek at HTML — if selector matches, use 'selector', else 'single'."""
    from playwright.async_api import async_playwright

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        await page.goto(html_path.resolve().as_uri(), wait_until="domcontentloaded")
        count = await page.evaluate(
            "sel => document.querySelectorAll(sel).length", selector
        )
        await browser.close()
    return "selector" if count > 0 else "single"


async def run(args: argparse.Namespace) -> int:
    html_path = Path(args.input).expanduser().resolve()
    if not html_path.exists():
        print(f"ERROR: input not found: {html_path}", file=sys.stderr)
        return 2

    width, height = args.size
    out_path = Path(args.output).expanduser().resolve()

    split = args.split
    if split == "auto":
        split = await detect_split_mode(html_path, args.selector)
        print(f"  • auto-detected split mode: {split}")

    with tempfile.TemporaryDirectory(prefix="html2pptx-") as tmp:
        tmp_dir = Path(tmp)
        if split == "selector":
            pngs = await render_via_selector(
                html_path, tmp_dir, args.selector, width, height, args.scale
            )
            if not pngs:
                print(
                    f"ERROR: --selector {args.selector!r} matched 0 elements. "
                    f"Try --split single or fix the selector.",
                    file=sys.stderr,
                )
                return 3
        elif split == "single":
            pngs = await render_single(html_path, tmp_dir, width, height, args.scale)
        else:
            print(f"ERROR: unknown --split {split!r}", file=sys.stderr)
            return 2

        print(f"  • rendered {len(pngs)} slide(s) at {width}x{height} @ {args.scale}x")
        assemble_pptx(pngs, out_path, width, height)

    print(f"✓ wrote {out_path}")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(
        description="Render an HTML file into a PPTX deck (one slide per page).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    ap.add_argument("--input", "-i", required=True, help="Path to source .html")
    ap.add_argument("--output", "-o", required=True, help="Path to output .pptx")
    ap.add_argument(
        "--size",
        type=parse_size,
        default=(1920, 1080),
        help="Slide pixel size WxH (default: 1920x1080)",
    )
    ap.add_argument(
        "--scale",
        type=float,
        default=2.0,
        help="Device pixel ratio for retina-quality screenshots (default: 2.0)",
    )
    ap.add_argument(
        "--split",
        choices=["auto", "selector", "single"],
        default="auto",
        help="How to split HTML into slides (default: auto)",
    )
    ap.add_argument(
        "--selector",
        default=DEFAULT_SELECTOR,
        help=f"CSS selector for slide elements (default: {DEFAULT_SELECTOR!r})",
    )
    args = ap.parse_args()
    return asyncio.run(run(args))


if __name__ == "__main__":
    sys.exit(main())
